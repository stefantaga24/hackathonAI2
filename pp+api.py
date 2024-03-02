import os
import random
from pptx import Presentation  # Import the PowerPoint library
from grazie.api.client.gateway import AuthType, GrazieApiGatewayClient, GrazieHeaders
from grazie.api.client.chat.prompt import ChatPrompt
from grazie.api.client.endpoints import GrazieApiGatewayUrls
from grazie.api.client.profiles import Profile


def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

token1 = "eyJhbGciOiJSUzUxMiIsInR5cCI6IkpXVCJ9.eyJzdWIiOiJHcmF6aWUgQXV0aGVudGljYXRpb24iLCJ1aWQiOiI2ZmQ2ZWI1OC05Njg1LTRkZTctYmMwNy00NmIxYmE2NTEzN2EiLCJ0eXBlIjoiY3VzdG9tIiwibGljZW5zZSI6ImN1c3RvbV9hcHA6OThiODk4MGMtZTk3My00MWYxLWIxZGEtYTExZTRiOGI5ZWVjIiwibGljZW5zZV90eXBlIjoiamV0YnJhaW5zLWFpLmFwcGxpY2F0aW9uLnN0YW5kYXJkIiwiZXhwIjoxNzEwNDQ2MTAwfQ.hVHEt-lMtlbdoOnm0W7gZ58av0BV4RjBS8mGFsNB3gJiFrMubVDdMYHftQ46XcqcglVpY1gymi4PJkGEtcwL1AJ8kYfxDjQqTE7YdMcauN3HXNoEHbhsexwyG4fGgG_5z36oqv_6f5BKZLUNWgfcn4JBEcxbxPoh8vmxwognd-tIJZQOylu6cSWcgWSSCJDrIUJQrs587BDrbHPZc0FPJoQtxaj6jf9tjf0IcL1xS9xCLrp5M7EhzP36O2eINYPJeDRHKPGOdHmb8WyxP5T6dopI9sF7SwZ-gcLhXCR9HpMTOX6jJagMALhP3POTkZ8gg9wbAnD4BhszJgpWqZCF4Q"


client_ip = "{}.{}.{}.{}".format(*[str(random.randint(0, 255)) for octet in range(4)])
client = GrazieApiGatewayClient(
    url= "https://api.app.stgn.grazie.aws.intellij.net/",
    grazie_jwt_token=token1,
    auth_type=AuthType.APPLICATION,
)


pptx_file_path = r"C:\Users\User\Downloads\1.IntroAI-intro.pptx"


pptx_text = extract_text_from_pptx(pptx_file_path)


chat_prompt = (
    ChatPrompt()
    .add_system("You are a helpful assistant.")
    .add_user(pptx_text)
)


response = client.chat(
    chat=chat_prompt,
    profile=Profile.OPENAI_GPT_4,
    headers={
        GrazieHeaders.ORIGINAL_USER_IP: client_ip,
    }
)

print(response.content)
